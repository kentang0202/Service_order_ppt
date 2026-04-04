import eel
import re
import datetime
import os
from pptx import Presentation

# '.' 代表當前資料夾
eel.init('.')

def move_slide(prs, old_index, new_index):
    """ 將投影片從 old_index 移動到 new_index """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

@eel.expose
def start_process(data):
    try:
        # 1. 檔案存在檢查 (使用 os 模組)
        template_path = '2026主日司會模板python.pptx'
        if not os.path.exists(template_path):
            return {
                "status": "❌ 找不到模板檔案：請確認 '2026主日司會模板python.pptx' 是否在程式旁。",
                "preview": ""
            }

        prs = Presentation(template_path)
        bible_raw = data.get('bible_raw', '').strip()
        proverbs_raw = data.get('proverbs_raw', '').strip()
        topic_str = data.get('topic', '').strip()

        # 2. 智慧型日期偵測 (優先從經文抓取，支援 yyyy/mm/dd 或 yyyy年mm月dd日)
        date_match = re.search(r'(\d{4})[年./](\d{1,2})[月./](\d{1,2})', bible_raw)

        if date_match:
            y, m, d = date_match.groups()
            date_str = f"{y}年{int(m)}月{int(d)}日"
            file_date = f"{y}{int(m):0>2}{int(d):0>2}"  # 補零格式供檔名使用
            date_display = f"📅 偵測日期：{date_str}"
        else:
            # 沒抓到則使用今天日期
            today = datetime.date.today()
            date_str = f"{today.year}年{today.month}月{today.day}日"
            file_date = today.strftime("%Y%m%d")
            # 增加 HTML 備註標籤
            date_display = f"📅 偵測日期：{date_str} <span class='auto-date'>(無輸入，使用當前日期)</span>"

        # 3. Regex 解析經文與箴言
        # 經文：抓取 〈〉 標題與內容
        bible_matches = re.findall(r'〈(.*?)〉\s*(.*?)(?=\n?〈|$)', bible_raw, re.S)
        # 箴言：抓取角色與內容 (支援多種破折號)
        prov_regex = r'(司會|會眾|全體)\s*[–—－-]\s*(.*?)(?=\n?(?:司會|會眾|全體)\s*[–—－-]|(?:\(阿們！.*?\))?$|$)'
        proverbs_matches = re.findall(prov_regex, proverbs_raw, re.S)

        if not bible_matches:
            return {"status": "❌ 辨識失敗：經文請確保包含 〈〉 括號標題。", "preview": ""}

        # --- 4. 開始填寫 PPT 邏輯 ---
        bible_titles = [m[0].strip() for m in bible_matches]

        # A. 更新日期與目錄 (第 2, 9, 10 頁)
        prs.slides[1].shapes.title.text = date_str  # Index 1
        prs.slides[8].shapes.title.text = date_str  # Index 8
        prs.slides[9].shapes.title.text = "\n".join(bible_titles)  # Index 9

        # B. 箴言填寫 (第 4-8 頁，Index 3-7)
        for i, (role, text) in enumerate(proverbs_matches):
            if i < 5:
                # 根據你的 debug，內容框 idx 為 0
                prs.slides[3 + i].shapes.title.text = text.strip()

        # C. 經文處理 (第 11 頁原地更新 + 後續 add_slide)
        # 第一則經文
        prs.slides[10].shapes.title.text = bible_titles[0]
        try:
            prs.slides[10].placeholders[1].text = bible_matches[0][1].strip()
        except: pass

        # 第二則之後的經文 (使用母片 Index 7)
        for i in range(1, len(bible_matches)):
            new_slide = prs.slides.add_slide(prs.slide_layouts[7])
            new_slide.shapes.title.text = bible_titles[i]
            try:
                new_slide.placeholders[1].text = bible_matches[i][1].strip()
            except: pass

        # D. 主題頁精準填寫 (第 12 頁，Index 11)
        topic_slide = prs.slides[11]
        titles_joined = "、".join(bible_titles)
        for shape in topic_slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = topic_str
            elif shape.placeholder_format.idx == 1:
                shape.text = titles_joined
            elif shape.placeholder_format.idx == 21:
                shape.text = f"{date_str} 主日禮拜"

        # E. 新建回顧頁與主題頁搬家
        review_slide = prs.slides.add_slide(prs.slide_layouts[1])
        review_slide.shapes.title.text = date_str
        
        # 將原本 Index 11 的主題頁搬到最後一頁
        move_slide(prs, 11, len(prs.slides) - 1)

        # 5. 存檔與預覽回傳
        output_name = f"{file_date}主日司會模板.pptx"
        prs.save(output_name)

        # 6. 建立詳細預覽 HTML (顯示全部內容) 
        preview_html = f"<div class='preview-header'>{date_display}</div>"

        # 處理主題的換行 (把純文字換行轉為 HTML 換行)
        formatted_topic = topic_str.strip().replace('\n', '<br>')

        # 組合主題：標籤後面加個 <br>，內容包在 .topic-content 裡
        # 如果 topic_str 為空，我們給一個預設提示或留白
        topic_display = (
            f"<div class='preview-header'>"
            f"🖊️ 話語主題：<br>"
            f"<span class='topic-content'>{formatted_topic if formatted_topic else '（尚未輸入）'}</span>"
            f"</div>"
        )
        preview_html += topic_display
        preview_html += "<hr>"

        
        # 顯示經文全部內容
        preview_html += "<b>📖 辨識到的經文：</b><br>"
        for title, content in bible_matches:
            # 1. 先處理內容的換行與空白，存入變數，避開 f-string 內部的反斜線
            formatted_content = content.strip().replace('\n', '<br>')
            
            # 2. 使用這個變數來建立 HTML 字串
            entry_html = (
                f"<div style='margin-bottom:15px; background:#fff; padding:10px; "
                f"border-radius:5px; border:1px solid #eee;'>"
                f"<b>【{title}】</b><br>{formatted_content}</div>"
            )
            preview_html += entry_html
        
        preview_html += "<hr><b>💡 辨識到的箴言：</b><br>"
        # 顯示箴言全部內容
        for i, (role, text) in enumerate(proverbs_matches):
            # 移除 [:40] 限制，顯示完整內容
            preview_html += f"<div style='margin-bottom:8px;'><b>{i+1}. {role}：</b> {text.strip()}</div>"

        return {
            "success": True,
            "status": f"✅ 成功生成：{output_name}",
            "preview": preview_html
        }

    except Exception as e:
        return {
            "success": False,
            "status": f"❌ 系統錯誤：{str(e)}", 
            "preview": ""
        }

# 啟動 Eel 視窗
eel.start('index.html', size=(800, 750))